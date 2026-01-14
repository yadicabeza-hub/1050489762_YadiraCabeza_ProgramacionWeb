from pptx import Presentation
from pptx.util import Inches, Pt
import textwrap

prs = Presentation()

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullets_slide(title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(18)

def add_code_slide(title, code):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.6)
    width = Inches(9)
    height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    wrapped = textwrap.fill(code, width=120)
    p = tf.paragraphs[0]
    p.font.name = 'Consolas'
    p.font.size = Pt(12)
    p.text = wrapped

# Cargar contenido desde slides_content.md sería ideal, pero para simplicidad lo embedemos aquí
add_title_slide('Sistema: Gestión de Productos', 'Laravel — CRUD Productos y Categorías')

add_bullets_slide('Agenda', [
    'Introducción y objetivo',
    'Stack y dependencias',
    'Arquitectura general',
    'Flujo CRUD (ejemplo)',
    'Demo y mejoras'
])

add_bullets_slide('Stack y dependencias', [
    'PHP ^8.2, Laravel ^12',
    'laravel/breeze (auth), phpunit',
    'Composer y npm para assets'
])

add_bullets_slide('Arquitectura general', [
    'Browser -> Rutas -> Middleware (auth) -> Controladores',
    'Modelos Eloquent -> Base de datos (migrations)',
    'Vistas Blade renderizando la UI'
])

add_bullets_slide('Flujo: Crear producto', [
    'Formulario -> POST /products con CSRF',
    'Validación en controller',
    'Product::create($validated)',
    'Redirección con mensaje de éxito'
])

code_example = '''$products = Product::with('category')->latest()->paginate(10);
return view('products.index', compact('products'));
'''
add_code_slide('Ejemplo de código (index)', code_example)

validation_example = '''$validated = $request->validate([
    'nombre' => 'required|string|max:255',
    'precio' => 'required|numeric|min:0',
    'stock' => 'required|integer|min:0',
    'estado' => 'required|in:activo,inactivo',
    'category_id' => 'nullable|exists:categories,id',
]);
'''
add_code_slide('Ejemplo de validación', validation_example)

add_bullets_slide('Limitaciones actuales', [
    'No hay manejo de subida/almacenamiento de imagen',
    'Validación en controlador (mover a FormRequest)',
    'Faltan tests y políticas de autorización'
])

add_bullets_slide('Pasos para demo local', [
    'composer install',
    'copiar .env.example -> .env y generar key',
    'php artisan migrate',
    'php artisan serve' 
])

add_bullets_slide('Mejoras recomendadas', [
    'Implementar Storage para imágenes',
    'Crear ProductRequest',
    'Añadir Feature tests',
    'Añadir Policies para autorización'
])

add_bullets_slide('Preguntas', ['¿Qué quieren ver con más detalle?'])

output_path = 'presentacion_proyecto_pw.pptx'
prs.save(output_path)
print(f'Presentación generada: {output_path}')
